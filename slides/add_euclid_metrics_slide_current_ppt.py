from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt


PPT_PATH = Path(__file__).with_name("audio-database-presentation.pptx")

INK = RGBColor(0x2A, 0x21, 0x19)
MUTED = RGBColor(0x72, 0x66, 0x59)
ACCENT = RGBColor(0x8C, 0x7B, 0x68)
SANS = "Aptos"
SERIF = "Georgia"
MONO = "Consolas"

EUCLID_TITLE = "Khoảng cách Euclid"
DETAIL_TITLE = "12 đại lượng dùng cho Euclid"


def set_text(shape, text, font_name, font_size, color, bold=False):
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold


def extract_title(slide):
    texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
    return texts[1] if len(texts) > 1 else ""


def duplicate_slide(prs, source_slide):
    new_slide = prs.slides.add_slide(prs.slide_layouts[6])
    new_slide.background.fill.solid()
    try:
      new_slide.background.fill.fore_color.rgb = source_slide.background.fill.fore_color.rgb
    except Exception:
      pass

    for shape in source_slide.shapes:
        new_element = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_element, "p:extLst")

    return new_slide


def move_slide(prs, old_index, new_index):
    slide_list = prs.slides._sldIdLst
    slide = slide_list[old_index]
    slide_list.remove(slide)
    slide_list.insert(new_index, slide)


def find_slide_index_by_title(prs, title):
    for index, slide in enumerate(prs.slides):
        if extract_title(slide) == title:
            return index
    return None


def update_euclid_slide(slide):
    set_text(slide.shapes[1], "CÔNG THỨC", SANS, 17, ACCENT, bold=True)
    set_text(slide.shapes[2], EUCLID_TITLE, SERIF, 34, INK)
    set_text(
        slide.shapes[4],
        (
            "Euclid(a, b) = sqrt(Σ(a_i - b_i)^2)\n"
            "Similarity = 1 / (1 + Euclid)\n\n"
            "Ký hiệu:\n"
            "- a, b: 2 vector so sánh 12 chiều\n"
            "- a_i, b_i: giá trị đã chuẩn hóa và áp trọng số ở chiều i\n\n"
            "Tham số trên FE:\n"
            "- result.similarity trên trang Search"
        ),
        MONO,
        13,
        INK,
    )
    set_text(
        slide.shapes[6],
        (
            "Trong code:\n"
            "distance = euclideanDistance(weightedSourceVector,\n"
            "  weightedCandidateVector);\n"
            "similarity = 1 / (1 + distance);"
        ),
        MONO,
        11,
        INK,
    )
    set_text(
        slide.shapes[7],
        (
            "Slide này đã được cập nhật theo phiên bản hiện tại của hệ thống: Euclid không còn so trực tiếp "
            "trên featureVector thô, mà so trên comparison profile 12 chiều rồi mới quy đổi ra phần trăm trùng khớp."
        ),
        SANS,
        14,
        MUTED,
    )


def update_detail_slide(slide):
    set_text(slide.shapes[1], "CÔNG THỨC", SANS, 17, ACCENT, bold=True)
    set_text(slide.shapes[2], DETAIL_TITLE, SERIF, 32, INK)
    set_text(
        slide.shapes[4],
        (
            "1. meanEnergy: năng lượng trung bình\n"
            "2. maxEnergy: năng lượng lớn nhất\n"
            "3. stdEnergy: độ lệch chuẩn năng lượng\n"
            "4. meanLoudness: độ to trung bình\n"
            "5. stdLoudness: độ lệch chuẩn loudness\n"
            "6. dominantPitchHz: cao độ trội"
        ),
        SANS,
        14,
        INK,
    )
    set_text(
        slide.shapes[6],
        (
            "7. normalizedBrightness: độ sáng chuẩn hóa\n"
            "8. meanZeroCrossingRate: tần suất đổi dấu TB\n"
            "9. maxPeak: biên độ đỉnh lớn nhất\n"
            "10. burstRatio = maxEnergy / meanEnergy\n"
            "11. activeWindowRatio: tỉ lệ cửa sổ hoạt động\n"
            "12. durationSeconds: thời lượng file"
        ),
        SANS,
        14,
        INK,
    )
    set_text(
        slide.shapes[7],
        (
            "Ánh xạ FE: kết quả cuối cùng hiển thị ở result.similarity. Nếu cần giải thích nguồn dữ liệu truy vấn, "
            "các đại lượng tóm tắt đi qua query.summary như dominantPitchHz, averageLoudnessDb, normalizedBrightness "
            "và durationSeconds."
        ),
        SANS,
        14,
        MUTED,
    )


def ensure_detail_slide(prs):
    existing_index = find_slide_index_by_title(prs, DETAIL_TITLE)
    euclid_index = find_slide_index_by_title(prs, EUCLID_TITLE)

    if euclid_index is None:
        raise RuntimeError("Không tìm thấy slide Euclid trong file PPT hiện tại.")

    euclid_slide = prs.slides[euclid_index]
    update_euclid_slide(euclid_slide)

    if existing_index is None:
        detail_slide = duplicate_slide(prs, euclid_slide)
        update_detail_slide(detail_slide)
        move_slide(prs, len(prs.slides) - 1, euclid_index + 1)
    else:
        update_detail_slide(prs.slides[existing_index])


def verify_formula_layout(prs):
    reference_index = find_slide_index_by_title(prs, EUCLID_TITLE)
    detail_index = find_slide_index_by_title(prs, DETAIL_TITLE)

    if reference_index is None or detail_index is None:
        return False, "Thiếu slide Euclid hoặc slide chi tiết mới."

    reference_slide = prs.slides[reference_index]
    detail_slide = prs.slides[detail_index]

    if len(reference_slide.shapes) != len(detail_slide.shapes):
        return False, "Số lượng shape khác với layout gốc."

    ref_signature = [
        (shape.shape_type, shape.left, shape.top, shape.width, shape.height)
        for shape in reference_slide.shapes
    ]
    detail_signature = [
        (shape.shape_type, shape.left, shape.top, shape.width, shape.height)
        for shape in detail_slide.shapes
    ]

    if ref_signature != detail_signature:
        return False, "Vị trí hoặc kích thước shape không còn khớp layout cũ."

    return True, "Layout vẫn khớp với slide công thức gốc."


def main():
    prs = Presentation(str(PPT_PATH))
    ensure_detail_slide(prs)
    prs.save(str(PPT_PATH))

    prs_check = Presentation(str(PPT_PATH))
    ok, message = verify_formula_layout(prs_check)
    print(message.encode("unicode_escape").decode())


if __name__ == "__main__":
    main()

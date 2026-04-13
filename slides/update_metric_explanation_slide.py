from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt


PPT_PATH = Path(__file__).with_name("audio-database-presentation.pptx")

INK = RGBColor(0x2A, 0x21, 0x19)
MUTED = RGBColor(0x72, 0x66, 0x59)
SERIF = "Georgia"
SANS = "Aptos"
MONO = "Consolas"


def set_text(shape, text, font_name, font_size, color, bold=False):
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold


def title_text(slide):
    texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
    return texts[1] if len(texts) > 1 else ""


def main():
    prs = Presentation(str(PPT_PATH))

    target_slide = None
    for slide in prs.slides:
        if title_text(slide) == "Ý nghĩa của hai chỉ số trên bảng tổng quan":
            target_slide = slide
            break

    if target_slide is None:
        raise RuntimeError("Không tìm thấy slide 'Giải thích chỉ số'.")

    # Shape mapping from current PPT structure:
    # 4: left title, 5: left body, 6: left note
    # 8: right title, 9: right body, 10: right note
    set_text(target_slide.shapes[4], "Số cửa sổ tín hiệu", SERIF, 24, INK)
    set_text(
        target_slide.shapes[5],
        "Đây là tổng số đoạn ngắn được tạo ra sau bước segmentation. "
        "Mỗi tệp âm thanh được chia thành nhiều cửa sổ nhỏ để hệ thống tính energy, "
        "loudness, pitch và brightness.\n\n"
        "Công thức chia segment:\n"
        "segmentCount = floor((sampleCount - windowSize) / hopSize) + 1",
        SANS,
        16,
        MUTED,
    )
    set_text(
        target_slide.shapes[6],
        "Ý nghĩa: số này càng lớn thì mức phân tích theo thời gian càng chi tiết.",
        SANS,
        16,
        MUTED,
    )

    set_text(target_slide.shapes[8], "Token metadata", SERIF, 24, INK)
    set_text(
        target_slide.shapes[9],
        "Đây là tổng số từ khóa mô tả được tách ra và đưa vào chỉ mục đảo metadata. "
        "Các token đến từ tiêu đề, mô tả, danh mục, thẻ và các trường metadata khác.\n\n"
        "Ví dụ:\n"
        "\"Tiếng chuông cảnh báo\" → [\"tiếng\", \"chuông\", \"cảnh\", \"báo\"]",
        SANS,
        16,
        MUTED,
    )
    set_text(
        target_slide.shapes[10],
        "Ý nghĩa: số này phản ánh quy mô chỉ mục metadata phục vụ tìm kiếm theo từ khóa.",
        SANS,
        16,
        MUTED,
    )

    prs.save(str(PPT_PATH))


if __name__ == "__main__":
    main()

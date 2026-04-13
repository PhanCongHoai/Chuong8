from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt


PPT_PATH = Path(__file__).with_name("audio-database-presentation.pptx")

INK = RGBColor(0x2A, 0x21, 0x19)
MUTED = RGBColor(0x72, 0x66, 0x59)
ACCENT = RGBColor(0x8C, 0x7B, 0x68)
SANS = "Aptos"
SERIF = "Georgia"
MONO = "Consolas"


FORMULA_SLIDES = [
    {
        "title": "Energy",
        "formula": "Energy = (1 / N) * Σ(x[n]^2)",
        "code": "Trong code:\nconst intensity = rms * rms;",
        "description": "Energy là mức năng lượng trung bình của cửa sổ tín hiệu.",
    },
    {
        "title": "Loudness",
        "formula": "Loudness = 20 * log10(RMS + ε)",
        "code": "Trong code:\nconst loudness = 20 * Math.log10(rms + EPSILON);",
        "description": "Loudness là độ to trên thang log của tín hiệu số.",
    },
    {
        "title": "Pitch",
        "formula": "Pitch ≈ sampleRate / bestLag",
        "code": "Trong code:\nreturn sampleRate / bestLag;",
        "description": "Pitch được ước lượng từ tự tương quan, lấy theo độ trễ mạnh nhất.",
    },
    {
        "title": "Brightness",
        "formula": "Brightness = Σ(f_k * |X[k]|) / Σ(|X[k]|)",
        "code": "Trong code:\nreturn weightedFrequency / magnitudeSum;",
        "description": "Brightness là trọng tâm phổ, phản ánh âm sắc sáng hay tối.",
    },
    {
        "title": "Khoảng cách Euclid",
        "formula": "Euclid(a, b) = sqrt(Σ(a_i - b_i)^2)\nĐiểm quy đổi = 1 / (1 + Euclid)",
        "code": (
            "Trong code:\n"
            "const difference = a[i] - b[i];\n"
            "squaredDistance += difference * difference;\n"
            "distance = Math.sqrt(squaredDistance);"
        ),
        "description": "Đây là công thức hệ thống đang dùng để xếp hạng âm thanh gần nhau.",
    },
    {
        "title": "Cosine Similarity",
        "formula": "Cosine(a, b) = (a · b) / (||a|| ||b||)",
        "code": (
            "Trong code tham khảo:\n"
            "dot += a[i] * b[i];\n"
            "magA += a[i] * a[i];\n"
            "magB += b[i] * b[i];"
        ),
        "description": "Công thức này được giữ trong slide để đối chiếu học thuật với Euclid.",
    },
]


def set_background_like(target_slide, source_slide):
    try:
        source_color = source_slide.background.fill.fore_color.rgb
    except Exception:
        return

    fill = target_slide.background.fill
    fill.solid()
    fill.fore_color.rgb = source_color


def duplicate_slide(prs, source_slide):
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)
    set_background_like(new_slide, source_slide)

    for shape in source_slide.shapes:
        new_element = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_element, "p:extLst")

    return new_slide


def move_slide(prs, old_index, new_index):
    slide_list = prs.slides._sldIdLst
    slide = slide_list[old_index]
    slide_list.remove(slide)
    slide_list.insert(new_index, slide)


def set_textbox_text(shape, text, font_name, font_size, color, bold=False):
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold


def format_formula_slide(slide, formula_data):
    # Shape index mapping comes from the current formula slide template.
    eyebrow = slide.shapes[1]
    title = slide.shapes[2]
    left_text = slide.shapes[4]
    right_text = slide.shapes[6]
    description = slide.shapes[7]

    set_textbox_text(eyebrow, "CÔNG THỨC", SANS, 17, ACCENT, bold=True)
    set_textbox_text(title, formula_data["title"], SERIF, 36, INK)
    set_textbox_text(left_text, formula_data["formula"], MONO, 20, INK)
    set_textbox_text(right_text, formula_data["code"], MONO, 16, INK)
    set_textbox_text(description, formula_data["description"], SANS, 20, MUTED)


def title_text(slide):
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()
            if text != "CÔNG THỨC":
                return text
    return ""


def main():
    prs = Presentation(str(PPT_PATH))

    source_index = None
    for index, slide in enumerate(prs.slides):
        if title_text(slide) == "Các công thức chính dùng trong đồ án":
            source_index = index
            break

    if source_index is None:
        raise RuntimeError("Không tìm thấy slide công thức gốc để tách.")

    source_slide = prs.slides[source_index]
    generated_slides = [source_slide]

    for _ in range(len(FORMULA_SLIDES) - 1):
        generated_slides.append(duplicate_slide(prs, source_slide))

    for slide, formula_data in zip(generated_slides, FORMULA_SLIDES):
        format_formula_slide(slide, formula_data)

    # Đưa các slide mới vừa tạo vào ngay sau slide công thức gốc.
    first_new_slide_index = len(prs.slides) - (len(FORMULA_SLIDES) - 1)
    insert_position = source_index + 1

    for current_index in range(first_new_slide_index, len(prs.slides)):
        move_slide(prs, current_index, insert_position)
        insert_position += 1

    prs.save(str(PPT_PATH))


if __name__ == "__main__":
    main()

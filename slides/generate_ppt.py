from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_FILE = Path(__file__).with_name("audio-database-presentation.pptx")

BG = RGBColor(0xF5, 0xF1, 0xE8)
PAPER = RGBColor(0xFB, 0xF8, 0xF2)
INK = RGBColor(0x2A, 0x21, 0x19)
MUTED = RGBColor(0x72, 0x66, 0x59)
LINE = RGBColor(0xD8, 0xCF, 0xC0)
ACCENT = RGBColor(0x8C, 0x7B, 0x68)

SERIF = "Georgia"
SANS = "Aptos"
MONO = "Consolas"


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def style_text_frame(text_frame):
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0


def add_topline(slide):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.75),
        Inches(0.55),
        Inches(1.15),
        Pt(1.2),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_eyebrow(slide, text):
    box = slide.shapes.add_textbox(Inches(0.78), Inches(0.72), Inches(3.4), Inches(0.4))
    frame = box.text_frame
    style_text_frame(frame)
    run = frame.paragraphs[0].add_run()
    run.text = text.upper()
    run.font.name = SANS
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = ACCENT


def add_title(slide, text, top=1.05, width=8.4):
    box = slide.shapes.add_textbox(Inches(0.78), Inches(top), Inches(width), Inches(2.05))
    frame = box.text_frame
    style_text_frame(frame)
    run = frame.paragraphs[0].add_run()
    run.text = text
    run.font.name = SERIF
    run.font.size = Pt(38)
    run.font.color.rgb = INK


def add_body(slide, text, left, top, width, height, size=18, color=MUTED):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    style_text_frame(frame)
    run = frame.paragraphs[0].add_run()
    run.text = text
    run.font.name = SANS
    run.font.size = Pt(size)
    run.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height, rounded=True):
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    card = slide.shapes.add_shape(
        shape_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = PAPER
    card.line.color.rgb = LINE
    return card


def add_card_title(slide, text, left, top, width):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.58))
    frame = box.text_frame
    style_text_frame(frame)
    run = frame.paragraphs[0].add_run()
    run.text = text
    run.font.name = SERIF
    run.font.size = Pt(22)
    run.font.color.rgb = INK


def add_bullet_list(slide, items, left, top, width, height, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    style_text_frame(frame)
    frame.clear()

    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.bullet = True
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(10)
        for run in paragraph.runs:
            run.font.name = SANS
            run.font.size = Pt(size)
            run.font.color.rgb = MUTED


def add_code_block(slide, code, left, top, width, height, size=15):
    add_card(slide, left, top, width, height)
    box = slide.shapes.add_textbox(
        Inches(left + 0.18),
        Inches(top + 0.18),
        Inches(width - 0.36),
        Inches(height - 0.36),
    )
    frame = box.text_frame
    style_text_frame(frame)
    run = frame.paragraphs[0].add_run()
    run.text = code
    run.font.name = MONO
    run.font.size = Pt(size)
    run.font.color.rgb = INK


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.78), Inches(6.72), Inches(11.2), Inches(0.5))
    frame = box.text_frame
    style_text_frame(frame)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = SANS
    run.font.size = Pt(14)
    run.font.color.rgb = MUTED


def make_slide(prs, eyebrow, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_topline(slide)
    add_eyebrow(slide, eyebrow)
    add_title(slide, title)
    return slide


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = make_slide(prs, "Bảo vệ đồ án", "Thiết kế và xây dựng hệ cơ sở dữ liệu âm thanh")
    add_body(
        slide,
        "Đồ án hiện thực bài 1 đến bài 5 của Chương 8: từ thiết kế hàm nguyên thủy, "
        "AudioIndex, AudioSQL đến hệ thống full-stack chạy với dữ liệu âm thanh thật.",
        0.78,
        2.05,
        6.6,
        1.5,
        20,
    )
    add_card(slide, 8.0, 1.35, 4.3, 4.7)
    add_card_title(slide, "Thông tin chính", 8.3, 1.75, 3.0)
    add_body(slide, "Môn học\nMultimedia Database Systems", 8.3, 2.2, 3.5, 0.9, 18)
    add_body(slide, "Công nghệ\nReact, Node.js, Express", 8.3, 3.2, 3.5, 0.9, 18)
    add_body(
        slide,
        "Dữ liệu thật\nTiếng vỗ tay, tiếng trống, tiếng chuông",
        8.3,
        4.2,
        3.5,
        1.0,
        18,
    )

    slide = make_slide(prs, "Đề bài", "Bài toán mà đồ án phải giải quyết")
    add_bullet_list(
        slide,
        [
            "Bài 1: Thiết kế tập hàm nguyên thủy cho audio retrieval.",
            "Bài 2: Xây cấu trúc AudioIndex mở rộng từ CreateAudioIndex.",
            "Bài 3: Thiết kế AudioSQL để truy vấn cơ sở dữ liệu âm thanh.",
            "Bài 4: Lập chỉ mục metadata cho âm thanh.",
            "Bài 5: Xây dựng hệ thống hoàn chỉnh hiện thực tất cả chức năng trên.",
        ],
        0.82,
        2.0,
        7.0,
        4.1,
        18,
    )
    add_card(slide, 8.25, 2.05, 4.05, 2.7)
    add_card_title(slide, "Tinh thần đề bài", 8.55, 2.35, 3.0)
    add_body(
        slide,
        "Âm thanh không chỉ được lưu như tệp, mà phải trở thành dữ liệu có thể "
        "phân tích, lập chỉ mục và truy vấn theo nội dung.",
        8.55,
        2.85,
        3.15,
        1.35,
        17,
    )

    slide = make_slide(prs, "Mục tiêu", "Bốn mục tiêu triển khai")
    goals = [
        ("Quản lý dữ liệu", "Lưu trữ âm thanh, metadata, bộ sưu tập và trạng thái phân tích."),
        ("Phân tích tín hiệu", "Chia đoạn, trích đặc trưng và tạo vector đại diện."),
        ("Tìm kiếm", "Tìm theo metadata hoặc tìm âm thanh tương tự."),
        ("Minh họa học thuật", "Thể hiện rõ quan hệ giữa bài 1, 2, 3, 4 và 5."),
    ]
    for index, (title, text) in enumerate(goals):
        left = 0.82 + index * 3.15
        add_card(slide, left, 2.0, 2.45, 3.1)
        add_card_title(slide, title, left + 0.18, 2.22, 2.0)
        add_body(slide, text, left + 0.18, 2.9, 2.0, 1.6, 15)

    slide = make_slide(prs, "Kiến trúc", "Kiến trúc ba lớp của đồ án")
    architecture = [
        ("Lớp 1", "Giao diện người dùng", "React hiển thị tổng quan, kho âm thanh, tìm kiếm và AudioSQL."),
        ("Lớp 2", "Dịch vụ nghiệp vụ", "Express cung cấp API cho tải lên, phân tích lại, truy vấn và xóa bản ghi."),
        ("Lớp 3", "Xử lý tín hiệu và chỉ mục", "Segmentation, feature extraction, VectorIndex và MetadataIndex."),
    ]
    for index, (label, title, text) in enumerate(architecture):
        top = 2.0 + index * 1.42
        add_card(slide, 0.82, top, 11.55, 1.1)
        add_body(slide, label, 1.05, top + 0.16, 0.9, 0.25, 14)
        add_card_title(slide, title, 2.0, top + 0.12, 3.0)
        add_body(slide, text, 5.0, top + 0.18, 6.3, 0.35, 15)

    slide = make_slide(prs, "Cơ sở lý thuyết", "Luồng xử lý theo Chương 8")
    steps = ["Đọc tín hiệu", "Cắt âm thanh", "Trích đặc trưng", "Tạo vector", "Lập chỉ mục", "Truy vấn"]
    for index, step in enumerate(steps):
        left = 0.82 + index * 2.03
        add_card(slide, left, 2.0, 1.72, 1.05)
        add_body(slide, f"{index + 1:02d}", left + 0.18, 2.14, 0.3, 0.2, 14)
        add_card_title(slide, step, left + 0.18, 2.42, 1.25)
    add_body(
        slide,
        "Âm thanh được xem là tín hiệu theo thời gian. Muốn truy vấn theo nội dung, "
        "hệ thống phải chia tín hiệu thành nhiều đoạn nhỏ rồi biến mỗi đoạn thành "
        "các đặc trưng có thể so sánh.",
        0.82,
        3.7,
        7.6,
        1.35,
        18,
    )
    add_card(slide, 8.75, 3.55, 3.55, 1.55)
    add_card_title(slide, "Ý tưởng cốt lõi", 9.0, 3.86, 2.2)
    add_body(
        slide,
        "Tệp âm thanh → nhiều cửa sổ nhỏ → vector đặc trưng → truy vấn gần nhất",
        9.0,
        4.28,
        2.9,
        0.55,
        17,
    )

    slide = make_slide(prs, "Cắt âm thanh", "Âm thanh được cắt như thế nào")
    add_bullet_list(
        slide,
        [
            "Tệp âm thanh được đưa về mono để xử lý thống nhất.",
            "Window size lấy theo khoảng 50 ms của tín hiệu.",
            "Hop size bằng một nửa window size, tương đương chồng lấn khoảng 50%.",
            "Mỗi cửa sổ tạo ra một bộ đặc trưng cục bộ.",
            "Toàn bộ đặc trưng cục bộ được lấy trung bình để tạo vector toàn cục.",
        ],
        0.82,
        2.0,
        5.35,
        4.2,
        17,
    )
    add_code_block(
        slide,
        "const windowSize = Math.max(1024, Math.floor(sampleRate * 0.05));\n"
        "const hopSize = Math.max(512, Math.floor(windowSize / 2));\n\n"
        "for (let start = 0; start + windowSize <= samples.length; start += hopSize) {\n"
        "  const segment = samples.slice(start, start + windowSize);\n"
        "  const features = computeWindowFeatures(segment, sampleRate, analysisWindow);\n"
        "  windows.push({ startSecond, endSecond, ...features });\n"
        "}",
        6.55,
        1.95,
        5.8,
        4.2,
    )
    add_footer(slide, "Đây chính là phần cắt âm thanh đã có trong code và đã được trình bày rõ trong slide.")

    slide = make_slide(prs, "Đặc trưng", "Giải thích các hàm đặc trưng âm thanh")
    features = [
        ("Energy", "Đo mức năng lượng trung bình của tín hiệu trong một cửa sổ."),
        ("Loudness", "Đưa biên độ về thang log để gần hơn với cảm nhận nghe."),
        ("Pitch", "Ước lượng cao độ từ độ trễ mạnh nhất trong autocorrelation."),
        ("Brightness", "Phản ánh độ sáng âm sắc qua trọng tâm phổ spectral centroid."),
        ("Zero-crossing rate", "Mật độ đổi dấu của tín hiệu trong một cửa sổ."),
        ("Peak", "Biên độ lớn nhất xuất hiện trong cửa sổ tín hiệu."),
    ]
    for index, (title, text) in enumerate(features):
        left = 0.82 + (index % 2) * 5.85
        top = 2.0 + (index // 2) * 1.42
        add_card(slide, left, top, 5.45, 1.15)
        add_card_title(slide, title, left + 0.18, top + 0.12, 2.5)
        add_body(slide, text, left + 2.22, top + 0.2, 2.9, 0.55, 15)

    slide = make_slide(prs, "Công thức", "Các công thức chính dùng trong đồ án")
    add_code_block(
        slide,
        "Energy = (1 / N) * Σ(x[n]^2)\n"
        "Loudness = 20 * log10(RMS + ε)\n"
        "Pitch ≈ sampleRate / bestLag\n"
        "Brightness = Σ(f_k * |X[k]|) / Σ(|X[k]|)",
        0.82,
        2.0,
        5.75,
        2.95,
    )
    add_code_block(
        slide,
        "Euclid(a, b) = sqrt(Σ(a_i - b_i)^2)\n"
        "Điểm quy đổi = 1 / (1 + Euclid)\n\n"
        "Cosine(a, b) = (a · b) / (||a|| ||b||)",
        6.63,
        2.0,
        5.72,
        2.95,
    )
    add_bullet_list(
        slide,
        [
            "Energy là mức năng lượng của cửa sổ tín hiệu.",
            "Loudness là mức độ to trên thang log của tín hiệu số.",
            "Pitch được lấy từ tự tương quan, không còn dùng zero-crossing rate làm công thức chính.",
            "Brightness là trọng tâm phổ, phản ánh âm sắc sáng hay tối.",
            "Euclid là công thức hệ thống đang dùng để xếp hạng âm thanh gần nhau.",
            "Cosine similarity vẫn được đưa vào slide để đối chiếu học thuật.",
        ],
        0.82,
        5.15,
        11.0,
        1.25,
        16,
    )

    slide = make_slide(prs, "Mã nguồn", "Đoạn mã tính đặc trưng và so sánh vector")
    add_code_block(
        slide,
        "const rms = computeRms(windowedSamples);\n"
        "const intensity = rms * rms;\n"
        "const loudness = 20 * Math.log10(rms + EPSILON);\n"
        "const pitch = estimatePitch(windowedSamples, sampleRate);\n"
        "const brightness = computeSpectralCentroid(windowedSamples, sampleRate);",
        0.82,
        2.0,
        5.75,
        3.25,
    )
    add_code_block(
        slide,
        "const difference = a[i] - b[i];\n"
        "squaredDistance += difference * difference;\n\n"
        "distance = Math.sqrt(squaredDistance);\n"
        "similarity = 1 / (1 + distance);",
        6.63,
        2.0,
        5.72,
        3.25,
    )
    add_footer(slide, "Đồ án hiện dùng khoảng cách Euclid để xếp hạng, đồng thời vẫn trình bày cosine similarity để đối chiếu.")

    slide = make_slide(prs, "Lập chỉ mục", "Hai loại chỉ mục của hệ thống")
    add_card(slide, 0.82, 2.0, 5.55, 2.15)
    add_card_title(slide, "Chỉ mục vector", 1.08, 2.28, 2.6)
    add_body(
        slide,
        "Mỗi bản ghi âm thanh có một vector toàn cục gồm các đặc trưng trung bình. "
        "Chỉ mục này phục vụ truy vấn âm thanh tương tự bằng khoảng cách Euclid.",
        1.08,
        2.78,
        4.65,
        1.1,
        17,
    )
    add_card(slide, 6.8, 2.0, 5.55, 2.15)
    add_card_title(slide, "Chỉ mục metadata", 7.06, 2.28, 2.7)
    add_body(
        slide,
        "Metadata được token hóa và đưa vào inverted index để tìm theo từ khóa, "
        "danh mục, thẻ và người phụ trách.",
        7.06,
        2.78,
        4.5,
        1.05,
        17,
    )
    add_code_block(
        slide,
        "for (const token of tokenize(searchableFields.join(\" \"))) {\n"
        "  addToBucket(tokenIndex, token, audio.id);\n"
        "}\n\n"
        "for (const tag of audio.tags || []) {\n"
        "  addToBucket(tagIndex, normalizeText(tag), audio.id);\n"
        "}",
        0.82,
        4.45,
        11.53,
        1.85,
    )

    slide = make_slide(prs, "Giải thích chỉ số", "Ý nghĩa của hai chỉ số trên bảng tổng quan")
    add_card(slide, 0.82, 2.0, 5.55, 3.2)
    add_card_title(slide, "Số cửa sổ tín hiệu", 1.08, 2.28, 3.0)
    add_body(
        slide,
        "Đây là tổng số đoạn ngắn được tạo ra sau bước segmentation. Mỗi tệp âm thanh "
        "được chia thành nhiều cửa sổ nhỏ để hệ thống tính energy, loudness, pitch và brightness.",
        1.08,
        2.8,
        4.6,
        1.3,
        17,
    )
    add_body(
        slide,
        "Ý nghĩa: số này càng lớn thì mức phân tích theo thời gian càng chi tiết.",
        1.08,
        4.35,
        4.55,
        0.35,
        16,
    )
    add_card(slide, 6.8, 2.0, 5.55, 3.2)
    add_card_title(slide, "Token metadata", 7.06, 2.28, 2.6)
    add_body(
        slide,
        "Đây là tổng số từ khóa mô tả được tách ra và đưa vào inverted index. "
        "Các token đến từ tiêu đề, mô tả, danh mục, thẻ và các trường metadata khác.",
        7.06,
        2.8,
        4.55,
        1.3,
        17,
    )
    add_body(
        slide,
        "Ý nghĩa: số này phản ánh quy mô chỉ mục metadata phục vụ tìm kiếm theo từ khóa.",
        7.06,
        4.35,
        4.55,
        0.35,
        16,
    )

    slide = make_slide(prs, "AudioSQL", "Cầu nối giữa bài 1, 2, 3 và 4")
    add_bullet_list(
        slide,
        [
            "AddAudioSource",
            "SegmentAudio",
            "ExtractAudioFeatures",
            "CreateAudioIndex",
            "SearchAudioByMetadata",
            "FindSimilarAudio",
            "ExecuteAudioSQL",
        ],
        0.82,
        2.0,
        4.6,
        4.0,
        17,
    )
    add_code_block(
        slide,
        "SELECT * FROM audios;\n"
        "SELECT * FROM audios WHERE SIMILAR_TO = 'audio-001' LIMIT 3;\n"
        "SHOW PRIMITIVE FUNCTIONS;\n"
        "SHOW AUDIO INDEX;",
        6.2,
        2.0,
        6.15,
        2.9,
    )
    add_body(
        slide,
        "AudioSQL là lớp truy vấn học thuật dùng để minh họa cách gọi các chức năng đã thiết kế ở bài 1 và bài 2.",
        6.25,
        5.1,
        5.6,
        0.9,
        17,
    )

    slide = make_slide(prs, "Dữ liệu thật", "Bộ dữ liệu dùng để demo")
    real_data = [
        ("Tiếng vỗ tay", "Âm thanh ngắn, đỉnh năng lượng cao và dễ nhận ra."),
        ("Tiếng trống", "Âm trầm hơn, biên độ lớn và có cấu trúc nhịp rõ hơn."),
        ("Tiếng chuông", "Âm sáng hơn, ngân dài và có phổ khác biệt."),
    ]
    for index, (title, text) in enumerate(real_data):
        left = 0.82 + index * 3.86
        add_card(slide, left, 2.0, 3.45, 2.25)
        add_card_title(slide, title, left + 0.2, 2.28, 2.35)
        add_body(slide, text, left + 0.2, 2.82, 2.75, 0.85, 15)
    add_body(
        slide,
        "Ba lớp âm thanh này đủ khác nhau để minh họa rõ việc trích đặc trưng, "
        "lập chỉ mục và truy vấn tương tự trên dữ liệu thật.",
        0.82,
        4.85,
        8.3,
        0.9,
        18,
    )

    slide = make_slide(prs, "Đánh giá", "Kết quả đạt được")
    add_bullet_list(
        slide,
        [
            "Hoàn thành đầy đủ yêu cầu bài 1 đến bài 5.",
            "Có hệ thống full-stack chạy được với dữ liệu âm thanh thật.",
            "Có phân tích tín hiệu, vector đặc trưng, metadata index và AudioSQL.",
            "Đã đổi phần so sánh vector sang khoảng cách Euclid theo yêu cầu.",
        ],
        0.82,
        2.0,
        6.25,
        4.0,
        18,
    )
    add_card(slide, 7.65, 2.05, 4.65, 2.6)
    add_card_title(slide, "Kết luận ngắn", 7.95, 2.35, 2.4)
    add_body(
        slide,
        "Đồ án đã biến tín hiệu âm thanh thô thành dữ liệu có cấu trúc, có chỉ mục "
        "và có thể truy vấn theo nội dung.",
        7.95,
        2.95,
        3.75,
        1.1,
        17,
    )

    slide = make_slide(prs, "Kết thúc", "Xin cảm ơn thầy cô đã lắng nghe")
    add_body(
        slide,
        "Hệ thống sẵn sàng cho phần demo trực tiếp: tải dữ liệu, xem đặc trưng tín hiệu, "
        "chạy AudioSQL và tìm âm thanh tương tự.",
        0.82,
        2.05,
        7.0,
        1.15,
        20,
    )
    add_card(slide, 8.05, 1.9, 4.25, 2.85)
    add_card_title(slide, "Sẵn sàng demo", 8.4, 2.28, 2.2)
    add_body(
        slide,
        "1. Tải tệp âm thanh\n2. Xem đặc trưng tín hiệu\n3. Tìm âm thanh tương tự\n4. Chạy AudioSQL",
        8.4,
        2.95,
        3.2,
        1.3,
        17,
    )

    prs.save(str(OUTPUT_FILE))


if __name__ == "__main__":
    build_presentation()

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt


PPT_PATH = Path(__file__).with_name("audio-database-presentation.pptx")

INK = RGBColor(0x2A, 0x21, 0x19)
MUTED = RGBColor(0x72, 0x66, 0x59)
ACCENT = RGBColor(0x8C, 0x7B, 0x68)
SANS = "Aptos"
SERIF = "Georgia"
MONO = "Consolas"


FORMULA_UPDATES = {
    "Energy": {
        "formula": (
            "Energy = (1 / N) * \u03a3(x[n]^2)\n\n"
            "K\u00fd hi\u1ec7u:\n"
            "- N: s\u1ed1 m\u1eabu trong 1 c\u1eeda s\u1ed5\n"
            "- x[n]: bi\u00ean \u0111\u1ed9 m\u1eabu th\u1ee9 n\n\n"
            "Tham s\u1ed1 tr\u00ean FE:\n"
            "- windows[i].energy\n"
            "- featureVector[0]"
        ),
        "code": (
            "Trong code:\n"
            "const intensity = rms * rms;\n"
            "energy: round(intensity);"
        ),
        "description": (
            "Energy l\u00e0 n\u0103ng l\u01b0\u1ee3ng trung b\u00ecnh c\u1ee7a 1 c\u1eeda s\u1ed5 t\u00edn hi\u1ec7u. "
            "Tr\u00ean frontend hi\u1ec7n t\u1ea1i ch\u01b0a c\u00f3 card ri\u00eang cho energy, nh\u01b0ng n\u00f3 v\u1eabn "
            "\u0111\u01b0\u1ee3c \u0111\u01b0a v\u00e0o feature vector t\u1ed5ng h\u1ee3p."
        ),
    },
    "Loudness": {
        "formula": (
            "Loudness = 20 * log10(RMS + \u03b5)\n"
            "RMS = sqrt((1 / N) * \u03a3(x[n]^2))\n\n"
            "K\u00fd hi\u1ec7u:\n"
            "- RMS: bi\u00ean \u0111\u1ed9 hi\u1ec7u d\u1ee5ng\n"
            "- \u03b5: s\u1ed1 nh\u1ecf \u0111\u1ec3 tr\u00e1nh log(0)\n\n"
            "Tham s\u1ed1 tr\u00ean FE:\n"
            "- windows[i].loudness\n"
            "- selectedAudio.summary.averageLoudnessDb"
        ),
        "code": (
            "Trong code:\n"
            "const loudness = 20 * Math.log10(rms + EPSILON);\n"
            "averageLoudnessDb = avg(windows[].loudness);"
        ),
        "description": (
            "Card '\u0110\u1ed9 to trung b\u00ecnh' tr\u00ean FE kh\u00f4ng l\u1ea5y gi\u00e1 tr\u1ecb c\u1ee7a 1 c\u1eeda s\u1ed5 "
            "\u0111\u01a1n l\u1ebb, m\u00e0 l\u00e0 gi\u00e1 tr\u1ecb loudness trung b\u00ecnh c\u1ee7a to\u00e0n b\u1ed9 c\u00e1c c\u1eeda s\u1ed5 trong t\u1ec7p."
        ),
    },
    "Pitch": {
        "formula": (
            "Pitch \u2248 sampleRate / bestLag\n\n"
            "K\u00fd hi\u1ec7u:\n"
            "- sampleRate: t\u1ea7n s\u1ed1 l\u1ea5y m\u1eabu (Hz)\n"
            "- bestLag: \u0111\u1ed9 tr\u1ec5 c\u00f3 t\u1ef1 t\u01b0\u01a1ng quan l\u1edbn nh\u1ea5t\n\n"
            "Tham s\u1ed1 tr\u00ean FE:\n"
            "- windows[i].pitch\n"
            "- selectedAudio.summary.dominantPitchHz"
        ),
        "code": (
            "Trong code:\n"
            "return sampleRate / bestLag;\n"
            "dominantPitchHz = avg(window.pitch > 0);"
        ),
        "description": (
            "Card 'Cao \u0111\u1ed9 tr\u1ed9i' tr\u00ean FE l\u00e0 trung b\u00ecnh pitch c\u1ee7a c\u00e1c c\u1eeda s\u1ed5 c\u00f3 pitch > 0, "
            "kh\u00f4ng ph\u1ea3i pitch t\u1ea1i m\u1ed9t th\u1eddi \u0111i\u1ec3m duy nh\u1ea5t."
        ),
    },
    "Brightness": {
        "formula": (
            "Brightness = \u03a3(f_k * |X[k]|) / \u03a3(|X[k]|)\n\n"
            "K\u00fd hi\u1ec7u:\n"
            "- f_k: t\u1ea7n s\u1ed1 c\u1ee7a bin ph\u1ed5 th\u1ee9 k\n"
            "- |X[k]|: \u0111\u1ed9 l\u1edbn ph\u1ed5 t\u1ea1i bin k\n\n"
            "Tham s\u1ed1 tr\u00ean FE:\n"
            "- windows[i].brightness\n"
            "- selectedAudio.summary.normalizedBrightness"
        ),
        "code": (
            "Trong code:\n"
            "return weightedFrequency / magnitudeSum;\n"
            "normalizedBrightness = avg(brightness) / (sampleRate / 2);"
        ),
        "description": (
            "Frontend hi\u1ec7n '\u0110\u1ed9 s\u00e1ng chu\u1ea9n h\u00f3a'. Ngh\u0129a l\u00e0 brightness trung b\u00ecnh "
            "\u0111\u00e3 \u0111\u01b0\u1ee3c chu\u1ea9n h\u00f3a theo t\u1ea7n s\u1ed1 Nyquist \u0111\u1ec3 \u0111\u01b0a v\u1ec1 kho\u1ea3ng [0, 1]."
        ),
    },
    "Kho\u1ea3ng c\u00e1ch Euclid": {
        "formula": (
            "Euclid(a, b) = sqrt(\u03a3(a_i - b_i)^2)\n"
            "Similarity = 1 / (1 + Euclid)\n\n"
            "K\u00fd hi\u1ec7u:\n"
            "- a_i, b_i: ph\u1ea7n t\u1eed th\u1ee9 i c\u1ee7a 2 vector \u0111\u1eb7c tr\u01b0ng\n\n"
            "Tham s\u1ed1 tr\u00ean FE:\n"
            "- result.similarity tren trang Search"
        ),
        "code": (
            "Trong code:\n"
            "const difference = a[i] - b[i];\n"
            "squaredDistance += difference * difference;\n"
            "distance = Math.sqrt(squaredDistance);\n"
            "similarity = 1 / (1 + distance);"
        ),
        "description": (
            "\u0110\u1ed9 t\u01b0\u01a1ng \u0111\u1ed3ng hi\u1ec7n tr\u00ean giao di\u1ec7n t\u00ecm ki\u1ebfm \u0111\u01b0\u1ee3c quy \u0111\u1ed5i t\u1eeb "
            "kho\u1ea3ng c\u00e1ch Euclid, kh\u00f4ng hi\u1ec7n tr\u1ef1c ti\u1ebfp gi\u00e1 tr\u1ecb distance."
        ),
    },
    "Cosine Similarity": {
        "formula": (
            "Cosine(a, b) = (a \u00b7 b) / (||a|| ||b||)\n\n"
            "K\u00fd hi\u1ec7u:\n"
            "- a \u00b7 b: t\u00edch v\u00f4 h\u01b0\u1edbng\n"
            "- ||a||, ||b||: \u0111\u1ed9 d\u00e0i vector\n\n"
            "Tham s\u1ed1 tr\u00ean FE:\n"
            "- Ch\u01b0a d\u00f9ng tr\u1ef1c ti\u1ebfp tr\u00ean FE"
        ),
        "code": (
            "Trong code tham kh\u1ea3o:\n"
            "dot += a[i] * b[i];\n"
            "magA += a[i] * a[i];\n"
            "magB += b[i] * b[i];"
        ),
        "description": (
            "Slide gi\u1eef c\u00f4ng th\u1ee9c cosine \u0111\u1ec3 \u0111\u1ed1i chi\u1ebfu h\u1ecdc thu\u1eadt. H\u1ec7 th\u1ed1ng hi\u1ec7n t\u1ea1i "
            "\u0111ang x\u1ebfp h\u1ea1ng b\u1eb1ng Euclid, n\u00ean frontend kh\u00f4ng hi\u1ec7n ch\u1ec9 s\u1ed1 cosine."
        ),
    },
}


def set_text(shape, text, font_name, font_size, color, bold=False):
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold


def extract_title(slide):
    texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
    return texts[1] if len(texts) > 1 else ""


def main():
    prs = Presentation(str(PPT_PATH))

    for slide in prs.slides:
        title = extract_title(slide)
        content = FORMULA_UPDATES.get(title)
        if not content:
            continue

        set_text(slide.shapes[1], "C\u00d4NG TH\u1ee8C", SANS, 17, ACCENT, bold=True)
        set_text(slide.shapes[2], title, SERIF, 34, INK)
        set_text(slide.shapes[4], content["formula"], MONO, 14, INK)
        set_text(slide.shapes[6], content["code"], MONO, 12, INK)
        set_text(slide.shapes[7], content["description"], SANS, 14, MUTED)

    prs.save(str(PPT_PATH))


if __name__ == "__main__":
    main()
